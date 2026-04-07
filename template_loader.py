#!/usr/bin/env python3
"""
Template Loader for PPTX Graphics Designer
Loads and applies PowerPoint templates for consistent styling
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pathlib import Path
import os

class TemplateLoader:
    """Loader for PowerPoint templates"""

    def __init__(self):
        self.templates = {}
        self._load_builtin_templates()

    def _load_builtin_templates(self):
        """Load built-in template configurations"""
        self.templates['reference'] = {
            'name': 'Reference Style',
            'colors': {
                'primary': RGBColor(66, 133, 244),
                'secondary': RGBColor(234, 67, 53),
                'accent': RGBColor(251, 188, 5),
                'background': RGBColor(255, 255, 255),
                'text': RGBColor(51, 51, 51),
            },
            'fonts': {
                'title': {'family': 'Microsoft YaHei', 'size': 32, 'bold': True},
                'body': {'family': 'Microsoft YaHei', 'size': 18, 'bold': False},
                'caption': {'family': 'Microsoft YaHei', 'size': 12, 'bold': False},
            },
            'layout': {
                'margin_left': 1.0,
                'margin_right': 1.0,
                'margin_top': 0.8,
                'margin_bottom': 0.8,
            }
        }

        self.templates['default'] = {
            'name': 'Default Style',
            'colors': {
                'primary': RGBColor(41, 128, 185),
                'secondary': RGBColor(39, 174, 96),
                'accent': RGBColor(142, 68, 173),
                'background': RGBColor(255, 255, 255),
                'text': RGBColor(51, 51, 51),
            },
            'fonts': {
                'title': {'family': 'Arial', 'size': 28, 'bold': True},
                'body': {'family': 'Arial', 'size': 16, 'bold': False},
                'caption': {'family': 'Arial', 'size': 11, 'bold': False},
            },
            'layout': {
                'margin_left': 1.0,
                'margin_right': 1.0,
                'margin_top': 1.0,
                'margin_bottom': 1.0,
            }
        }

        self.templates['custom'] = {
            'name': 'Custom Style',
            'colors': {
                'primary': RGBColor(26, 35, 126),
                'secondary': RGBColor(13, 71, 161),
                'accent': RGBColor(57, 73, 171),
                'background': RGBColor(255, 255, 255),
                'text': RGBColor(51, 51, 51),
            },
            'fonts': {
                'title': {'family': 'Calibri', 'size': 30, 'bold': True},
                'body': {'family': 'Calibri', 'size': 18, 'bold': False},
                'caption': {'family': 'Calibri', 'size': 12, 'bold': False},
            },
            'layout': {
                'margin_left': 0.8,
                'margin_right': 0.8,
                'margin_top': 0.8,
                'margin_bottom': 0.8,
            }
        }

    def load_template(self, template_path):
        """
        Load template from file path or built-in name

        Args:
            template_path: Path to .pptx template file or built-in template name

        Returns:
            dict: Template configuration
        """
        # Check if it's a built-in template
        if template_path in self.templates:
            return self.templates[template_path]

        # Check if it's a file path
        if os.path.isfile(template_path) and template_path.endswith('.pptx'):
            return self._load_from_pptx(template_path)

        # Default fallback
        return self.templates['default']

    def _load_from_pptx(self, pptx_path):
        """
        Load template configuration from PPTX file

        Args:
            pptx_path: Path to PPTX template file

        Returns:
            dict: Extracted template configuration
        """
        try:
            prs = Presentation(pptx_path)

            # Extract theme colors from slide master
            theme_colors = self._extract_theme_colors(prs)

            # Extract font information
            fonts = self._extract_fonts(prs)

            # Create template config
            template = {
                'name': Path(pptx_path).stem,
                'colors': theme_colors,
                'fonts': fonts,
                'layout': {
                    'margin_left': 1.0,
                    'margin_right': 1.0,
                    'margin_top': 0.8,
                    'margin_bottom': 0.8,
                },
                'source_file': pptx_path
            }

            return template

        except Exception as e:
            print(f"Warning: Could not load template from {pptx_path}: {e}")
            return self.templates['default']

    def _extract_theme_colors(self, prs):
        """Extract color scheme from presentation"""
        # Default colors
        colors = {
            'primary': RGBColor(41, 128, 185),
            'secondary': RGBColor(39, 174, 96),
            'accent': RGBColor(142, 68, 173),
            'background': RGBColor(255, 255, 255),
            'text': RGBColor(51, 51, 51),
        }

        try:
            # Try to extract from slide master
            if prs.slide_masters:
                master = prs.slide_masters[0]
                # This is a simplified extraction - real implementation would
                # parse the theme XML for actual colors
                pass
        except:
            pass

        return colors

    def _extract_fonts(self, prs):
        """Extract font information from presentation"""
        fonts = {
            'title': {'family': 'Arial', 'size': 28, 'bold': True},
            'body': {'family': 'Arial', 'size': 16, 'bold': False},
            'caption': {'family': 'Arial', 'size': 11, 'bold': False},
        }

        try:
            # Try to extract from slide layouts
            if prs.slide_layouts:
                layout = prs.slide_layouts[0]
                # Simplified font extraction
                pass
        except:
            pass

        return fonts

    def apply_template_to_slide(self, slide, template_config):
        """
        Apply template configuration to a slide

        Args:
            slide: PowerPoint slide object
            template_config: Template configuration dict
        """
        # This would apply the template settings to the slide
        # For now, it's a placeholder for future implementation
        pass

    def get_available_templates(self):
        """Get list of available built-in templates"""
        return list(self.templates.keys())

    def create_slide_from_template(self, prs, template_config, layout_index=5):
        """
        Create a new slide using template configuration

        Args:
            prs: Presentation object
            template_config: Template configuration
            layout_index: Slide layout index

        Returns:
            Slide object
        """
        slide = prs.slides.add_slide(prs.slide_layouts[layout_index])

        # Apply template settings
        self.apply_template_to_slide(slide, template_config)

        return slide

def load_template(template_path):
    """Convenience function to load a template"""
    loader = TemplateLoader()
    return loader.load_template(template_path)

def get_available_templates():
    """Get list of available templates"""
    loader = TemplateLoader()
    return loader.get_available_templates()
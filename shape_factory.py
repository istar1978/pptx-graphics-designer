#!/usr/bin/env python3
"""
Shape Factory for PPTX Graphics Designer
Provides standardized shape creation with advanced features
"""

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN

class ShapeFactory:
    """Factory for creating standardized shapes with advanced features"""

    def __init__(self, slide, colors):
        self.slide = slide
        self.colors = colors

    def create_box(self, x, y, w, h, text, subtitle=None, color='primary',
                   shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
                   font_size=11, subtitle_size=9,
                   dashed_border=False, border_width=2):
        """Create a box with main text and optional subtitle"""
        shape = self.slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))

        # Fill
        shape.fill.solid()
        fill_color = self.colors.get(color, self.colors.get('primary', RGBColor(41, 128, 185)))
        shape.fill.fore_color.rgb = fill_color

        # Border
        shape.line.color.rgb = self.colors.get('border', RGBColor(150, 150, 150))
        shape.line.width = Pt(border_width)
        if dashed_border:
            shape.line.dash_style = 4  # MSO_DASH_STYLE.DASH

        # Text frame
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        tf.clear()

        # Main text
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = True
        p.font.color.rgb = self.colors.get('text_dark', RGBColor(51, 51, 51))
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        if subtitle:
            p = tf.add_paragraph()
            p.text = subtitle
            p.font.size = Pt(subtitle_size)
            p.font.bold = False
            p.font.color.rgb = self.colors.get('text_light', RGBColor(102, 102, 102))
            p.alignment = PP_ALIGN.CENTER

        return shape

    def create_connector(self, begin_shape, end_shape, begin_side=1, end_side=3,
                        color='border', width=2, dashed=False):
        """Create a connector line between two shapes"""
        connector = self.slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(0), Inches(0), Inches(1), Inches(1)
        )

        # Set connection points
        connector.begin_connect(begin_shape, begin_side)
        connector.end_connect(end_shape, end_side)

        # Style
        connector.line.color.rgb = self.colors.get(color, RGBColor(150, 150, 150))
        connector.line.width = Pt(width)
        if dashed:
            connector.line.dash_style = 4

        return connector

    def create_textbox(self, x, y, w, h, text, font_size=12, alignment=PP_ALIGN.LEFT):
        """Create a simple textbox"""
        textbox = self.slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = textbox.text_frame
        tf.text = text
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)

        p = tf.paragraphs[0]
        p.font.size = Pt(font_size)
        p.font.color.rgb = self.colors.get('text_dark', RGBColor(51, 51, 51))
        p.alignment = alignment

        return textbox

    def create_circle(self, x, y, diameter, text, color='primary'):
        """Create a circle shape"""
        shape = self.slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(diameter), Inches(diameter)
        )

        # Fill
        shape.fill.solid()
        fill_color = self.colors.get(color, self.colors.get('primary', RGBColor(41, 128, 185)))
        shape.fill.fore_color.rgb = fill_color

        # Border
        shape.line.color.rgb = self.colors.get('border', RGBColor(150, 150, 150))
        shape.line.width = Pt(2)

        # Text
        tf = shape.text_frame
        tf.text = text
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.1)
        tf.margin_bottom = Inches(0.1)

        p = tf.paragraphs[0]
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = self.colors.get('white', RGBColor(255, 255, 255))
        p.alignment = PP_ALIGN.CENTER

        return shape

    def create_rectangle(self, x, y, w, h, text, color='primary', rounded=True):
        """Create a rectangle shape"""
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
        shape = self.slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))

        # Fill
        shape.fill.solid()
        fill_color = self.colors.get(color, self.colors.get('primary', RGBColor(41, 128, 185)))
        shape.fill.fore_color.rgb = fill_color

        # Border
        shape.line.color.rgb = self.colors.get('border', RGBColor(150, 150, 150))
        shape.line.width = Pt(2)

        # Text
        tf = shape.text_frame
        tf.text = text
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.1)
        tf.margin_bottom = Inches(0.1)

        p = tf.paragraphs[0]
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = self.colors.get('text_dark', RGBColor(51, 51, 51))
        p.alignment = PP_ALIGN.CENTER

        return shape
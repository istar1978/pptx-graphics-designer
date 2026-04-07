#!/usr/bin/env python3
"""
PPTX Graphics Designer - Professional diagram creation for PowerPoint
Supports: flowcharts, architecture diagrams, network topology, Gantt charts, sequence diagrams, scatter plots

Enhanced version with:
- Complex multi-layer architecture diagrams
- Multi-line text with subtitle support
- Gradient-like color fills
- Dashed borders and connectors
- Logo and header support
- Side panel for descriptions
"""

import argparse
import json
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.line import LineFormat

# Import our custom modules
from style_presets import get_color_scheme, COLOR_SCHEMES
from shape_factory import ShapeFactory
from layout_engine import LayoutEngine
from mermaid_parser import parse_mermaid
from template_loader import load_template

# ============================================================================
# Enhanced Shape Factory
# ============================================================================

# ============================================================================
# Advanced Layout Engine
# ============================================================================

# ============================================================================
# Advanced Diagram Builders
# ============================================================================

class EnterpriseArchitectureBuilder:
    """Build enterprise-style architecture diagrams (参考图风格)"""
    
    def __init__(self, slide, colors):
        self.slide = slide
        self.colors = colors
        self.factory = ShapeFactory(slide, colors)
    
    def build(self, layers_config, side_panel_items=None):
        """
        Build enterprise architecture diagram
        
        layers_config: [
            {
                'name': '接入层',
                'color': 'access',
                'bg_color': 'bg_access',
                'items': [
                    {'text': '400 热线', 'subtitle': None},
                    {'text': '官网客服', 'subtitle': 'Web'},
                    {'text': '意图识别', 'subtitle': 'BERT · Transformer'},
                ]
            },
            ...
        ]
        """
        slide_w = 13.333
        slide_h = 7.5
        
        # Calculate layer positions
        if side_panel_items:
            content_w = slide_w - 3.5  # Leave space for side panel
        else:
            content_w = slide_w - 1.0
        
        layer_h = 1.0
        gap = 0.12
        start_y = 1.0
        margin_x = 0.3
        
        layer_positions = []
        current_y = start_y
        for layer in layers_config:
            layer_positions.append((margin_x, current_y, content_w, layer_h))
            current_y += layer_h + gap
        
        # Draw layers
        for i, (layer, pos) in enumerate(zip(layers_config, layer_positions)):
            x, y, w, h = pos
            
            # Get colors
            bg_color = self.colors.get(layer.get('bg_color', 'white'))
            border_color = self.colors.get(layer.get('color', 'border'))
            
            # Create layer container
            self.factory.create_layer_container(
                x, y, w, h,
                layer['name'],
                bg_color, border_color,
                dashed=layer.get('dashed', False)
            )
            
            # Calculate item positions
            items = layer.get('items', [])
            num_items = len(items)
            if num_items > 0:
                item_w = (w - 0.6 - (num_items - 1) * 0.15) / num_items
                item_h = 0.65
                item_start_x = x + 0.3
                item_y = y + 0.45
                
                # Draw items
                for j, item_data in enumerate(items):
                    item_x = item_start_x + j * (item_w + 0.15)
                    
                    # Determine item color based on layer
                    item_color = self.colors.get(layer.get('color', 'primary'))
                    # Use lighter fill for items
                    item_bg = self.colors.get(layer.get('bg_color', 'white'))
                    
                    # Create item box
                    self.factory.create_box(
                        item_x, item_y, item_w, item_h,
                        item_data.get('text', f'Item {j+1}'),
                        item_data.get('subtitle'),
                        color=item_bg,
                        border_width=1.5,
                        font_size=10,
                        subtitle_size=8
                    )
        
        # Draw side panel if provided
        if side_panel_items:
            panel_x = slide_w - 3.3
            panel_y = 1.0
            panel_w = 3.0
            panel_h = slide_h - 1.5
            
            # Panel background
            panel_bg = self.slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(panel_x), Inches(panel_y),
                Inches(panel_w), Inches(panel_h)
            )
            panel_bg.fill.solid()
            panel_bg.fill.fore_color.rgb = RGBColor(245, 250, 255)
            panel_bg.line.fill.background()
            
            # Add items
            for section in side_panel_items:
                self.factory.create_side_panel(
                    panel_x + 0.2, panel_y + 0.2, panel_w - 0.4, 2.5,
                    section['title'],
                    section['items']
                )
                panel_y += 1.8

# ============================================================================
# Slide Creation Helpers
# ============================================================================

def create_slide(prs, title="", subtitle="", show_header=True, logo_path=None):
    """Create a slide with optional header and logo"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.fill.background()
    
    # Header
    if show_header:
        header_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8)
        )
        header_bg.fill.solid()
        header_bg.fill.fore_color.rgb = RGBColor(250, 252, 255)
        header_bg.line.fill.background()
        
        # Title
        if title:
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.2), Inches(10), Inches(0.5)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = RGBColor(51, 51, 51)
            
            # Decorative line
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.72), Inches(2), Inches(0.08)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(66, 133, 244)
            line.line.fill.background()
        
        # Subtitle
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(3.0), Inches(0.25), Inches(8), Inches(0.4)
            )
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(102, 102, 102)
    
    return slide

# ============================================================================
# Main Generation Function
# ============================================================================

def generate_enterprise_architecture(output_path, layers, side_panel=None, 
                                     title="技术架构", style='enterprise', 
                                     template_config=None, layout='auto'):
    """Generate enterprise-style architecture diagram"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    colors = get_color_scheme(style)
    slide = create_slide(prs, title, "", show_header=True)
    
    builder = EnterpriseArchitectureBuilder(slide, colors)
    builder.build(layers, side_panel)
    
    prs.save(output_path)
    print(f"✓ Enterprise architecture diagram generated: {output_path}")
    print(f"  Layers: {len(layers)}")
    print(f"  Style: {style}")

def generate_diagram(output_path, diagram_type, input_data, style='default', title='', 
                   template_config=None, layout='auto'):
    """Generate a diagram"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    colors = get_color_scheme(style)
    slide = create_slide(prs, title, "", show_header=False)
    
    # Simple architecture for backward compatibility
    if diagram_type == 'architecture':
        if isinstance(input_data, str):
            # Parse simple format: "Layer1: item1,item2|Layer2: item3,item4"
            layers = []
            for layer_str in input_data.split('|'):
                if ':' in layer_str:
                    name, items = layer_str.split(':', 1)
                    layers.append({
                        'name': name.strip(),
                        'items': [{'text': i.strip()} for i in items.split(',')]
                    })
        else:
            layers = input_data
        
        builder = EnterpriseArchitectureBuilder(slide, colors)
        builder.build(layers)
    
    prs.save(output_path)
    print(f"✓ Diagram generated: {output_path}")

# ============================================================================
# CLI
# ============================================================================

def main():
    parser = argparse.ArgumentParser(description='PPTX Graphics Designer v2.0')
    parser.add_argument('--input', '-i', required=True, help='Input JSON file or description')
    parser.add_argument('--output', '-o', required=True, help='Output PPTX file')
    parser.add_argument('--type', '-t', default='architecture',
                       choices=['flowchart', 'architecture', 'enterprise', 'network', 'gantt', 'sequence', 'scatter', 'mermaid'],
                       help='Diagram type')
    parser.add_argument('--style', default='enterprise',
                       choices=['enterprise', 'default', 'professional', 'colorful', 'minimal'],
                       help='Visual style')
    parser.add_argument('--title', default='技术架构', help='Slide title')
    parser.add_argument('--template', default='default',
                       help='PPT template file path or built-in template name')
    parser.add_argument('--layout', default='auto',
                       choices=['auto', 'vertical', 'horizontal', 'grid', 'composite'],
                       help='Layout arrangement')
    parser.add_argument('--side-panel', help='Side panel content (JSON or file path)')
    
    args = parser.parse_args()
    
    # Load input
    if args.input.endswith('.json'):
        with open(args.input, 'r', encoding='utf-8') as f:
            input_data = json.load(f)
    else:
        try:
            input_data = json.loads(args.input)
        except json.JSONDecodeError:
            input_data = args.input
    
    # Load side panel
    side_panel = None
    if args.side_panel:
        if args.side_panel.endswith('.json'):
            with open(args.side_panel, 'r', encoding='utf-8') as f:
                side_panel = json.load(f)
        else:
            try:
                side_panel = json.loads(args.side_panel)
            except json.JSONDecodeError:
                pass
    
    # Load template
    template_config = load_template(args.template)
    
    # Generate diagram based on type
    if args.type == 'mermaid':
        generate_mermaid_diagram(args.output, input_data, args.title, args.style, template_config, args.layout)
    elif args.type == 'enterprise' or isinstance(input_data, list):
        generate_enterprise_architecture(
            args.output, input_data, side_panel, args.title, args.style, template_config, args.layout
        )
    else:
        generate_diagram(args.output, args.type, input_data, args.style, args.title, template_config, args.layout)

# ============================================================================
# Mermaid Diagram Generation
# ============================================================================

def generate_mermaid_diagram(output_path, mermaid_code, title, style, template_config, layout):
    """Generate diagram from Mermaid syntax"""
    try:
        # Preprocess Mermaid code: convert single-line semicolon-separated input into lines
        if isinstance(mermaid_code, str) and ';' in mermaid_code and '\n' not in mermaid_code:
            mermaid_code = "\n".join([part.strip() for part in mermaid_code.split(';') if part.strip()])

        # Parse Mermaid code
        diagram_data = parse_mermaid(mermaid_code)

        if diagram_data is None or not isinstance(diagram_data, dict) or 'type' not in diagram_data:
            print("Mermaid parser returned invalid result. Falling back to text output.")
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11), Inches(5))
            tf = textbox.text_frame
            tf.text = f"Mermaid Parsing Failed\n\n{mermaid_code}"
            prs.save(output_path)
            return
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Add slide
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.5), Inches(0.7))
        title_tf = title_box.text_frame
        title_tf.text = title
        title_tf.paragraphs[0].font.size = Pt(24)
        title_tf.paragraphs[0].font.bold = True
        
        # Get colors
        colors = get_color_scheme(style)
        
        # Create layout engine
        layout_engine = LayoutEngine()
        
        # Generate diagram based on type
        if diagram_data['type'] == 'flowchart':
            _generate_flowchart_diagram(slide, diagram_data, colors, layout_engine)
        elif diagram_data['type'] == 'sequence':
            _generate_sequence_diagram(slide, diagram_data, colors, layout_engine)
        else:
            # Fallback to simple text representation
            _generate_text_diagram(slide, diagram_data, colors)
        
        # Save presentation
        prs.save(output_path)
        print(f"✓ Mermaid diagram generated: {output_path}")
        
    except Exception as e:
        print(f"Error generating Mermaid diagram: {e}")
        # Fallback to basic presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11), Inches(5))
        tf = textbox.text_frame
        tf.text = f"Mermaid Diagram\n\n{mermaid_code}"
        prs.save(output_path)

def _generate_flowchart_diagram(slide, diagram_data, colors, layout_engine):
    """Generate flowchart from parsed data"""
    factory = ShapeFactory(slide, colors)
    
    # Create nodes
    node_positions = {}
    nodes = diagram_data['nodes']
    
    # Simple positioning
    y_pos = 2.0
    for i, node in enumerate(nodes):
        x_pos = 2.0 + (i % 3) * 3.0
        if i % 3 == 0:
            y_pos += 1.5
        
        shape = factory.create_box(x_pos, y_pos, 2.0, 0.8, node['label'])
        node_positions[node['id']] = (x_pos + 1.0, y_pos + 0.4)  # Center point
    
    # Create edges
    for edge in diagram_data['edges']:
        source_pos = node_positions.get(edge['source'])
        target_pos = node_positions.get(edge['target'])

        if source_pos and target_pos:
            x1, y1 = source_pos
            x2, y2 = target_pos
            # Draw a straight connector between centers
            try:
                line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
                line.line.color.rgb = colors.get('border', RGBColor(150, 150, 150))
                line.line.width = Pt(2)
                if edge.get('style') == 'dashed':
                    line.line.dash_style = 4
            except Exception:
                # Fallback to shape-based line if connector isn't supported
                line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(min(x1, x2)), Inches(min(y1, y2)), Inches(abs(x2 - x1)), Pt(2))

            # Optional label
            if edge.get('label'):
                mid_x = (x1 + x2) / 2
                mid_y = (y1 + y2) / 2
                lbl = slide.shapes.add_textbox(Inches(mid_x - 0.6), Inches(mid_y - 0.2), Inches(1.2), Inches(0.4))
                tf = lbl.text_frame
                tf.text = edge.get('label')
                tf.paragraphs[0].font.size = Pt(9)
                tf.paragraphs[0].font.color.rgb = colors.get('text_light', RGBColor(102, 102, 102))

def _generate_sequence_diagram(slide, diagram_data, colors, layout_engine):
    """Generate sequence diagram from parsed data"""
    factory = ShapeFactory(slide, colors)
    
    # Create participants
    participants = diagram_data['participants']
    participant_positions = {}
    
    x_start = 1.5
    x_spacing = 2.5
    
    for i, participant in enumerate(participants):
        x_pos = x_start + i * x_spacing
        factory.create_box(x_pos, 1.0, 1.8, 0.6, participant['label'])
        participant_positions[participant['id']] = x_pos + 0.9  # Center
    
    # Create lifelines
    for i, participant in enumerate(participants):
        x_pos = x_start + i * x_spacing + 0.9
        # Draw vertical line (simplified)
        line = slide.shapes.add_shape(MSO_SHAPE.LINE, Inches(x_pos), Inches(1.6), Inches(0), Inches(4.0))
        line.line.color.rgb = colors.get('border', RGBColor(150, 150, 150))
    
    # Add messages
    y_pos = 2.0
    for message in diagram_data['messages']:
        source_x = participant_positions.get(message['source'])
        target_x = participant_positions.get(message['target'])
        
        if source_x and target_x:
            # Draw arrow
            arrow = slide.shapes.add_shape(MSO_SHAPE.LINE, Inches(source_x), Inches(y_pos), Inches(target_x), Inches(y_pos))
            arrow.line.color.rgb = colors.get('primary', RGBColor(41, 128, 185))
            
            # Add label
            label_x = (source_x + target_x) / 2
            factory.create_textbox(label_x - 0.5, y_pos - 0.2, 1.0, 0.4, message['message'], 8)
            
            y_pos += 0.8

def _generate_text_diagram(slide, diagram_data, colors):
    """Fallback text representation"""
    factory = ShapeFactory(slide, colors)
    
    content = f"Diagram Type: {diagram_data['type']}\n\n"
    
    if 'nodes' in diagram_data:
        content += "Nodes:\n"
        for node in diagram_data['nodes']:
            content += f"- {node.get('label', node.get('id', 'Unknown'))}\n"
    
    if 'edges' in diagram_data:
        content += "\nConnections:\n"
        for edge in diagram_data['edges']:
            content += f"- {edge['source']} -> {edge['target']}\n"
    
    factory.create_textbox(1.0, 1.5, 11.0, 5.0, content, 12)


if __name__ == '__main__':
    main()

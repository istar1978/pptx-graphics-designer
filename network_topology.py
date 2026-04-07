#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Network Topology Diagram Generator for PowerPoint
Supports: routers, switches, firewalls, servers, PCs, cloud, and network zones

Creates Visio-style network topology diagrams with:
- Standard network device icons (router, switch, firewall, server, PC, etc.)
- Network zone containers (DMZ, Office, Data Center, etc.)
- Connection lines with labels
- Professional Cisco-style iconography
"""

import argparse
import json
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN

# ============================================================================
# Network Device Icon Specifications
# ============================================================================

NETWORK_ICONS = {
    # Network Infrastructure
    'router': {
        'shape': MSO_SHAPE.OVAL,
        'color': RGBColor(0, 153, 204),
        'icon_type': 'router',
        'width_ratio': 1.2,
        'height_ratio': 0.8,
    },
    'switch': {
        'shape': MSO_SHAPE.OVAL,
        'color': RGBColor(0, 153, 204),
        'icon_type': 'switch',
        'width_ratio': 1.2,
        'height_ratio': 0.8,
    },
    'firewall': {
        'shape': MSO_SHAPE.CUBE,
        'color': RGBColor(255, 102, 0),
        'icon_type': 'firewall',
        'width_ratio': 1.0,
        'height_ratio': 1.0,
    },
    'load_balancer': {
        'shape': MSO_SHAPE.HEXAGON,
        'color': RGBColor(153, 0, 204),
        'icon_type': 'lb',
        'width_ratio': 1.1,
        'height_ratio': 0.9,
    },
    
    # Servers
    'server': {
        'shape': MSO_SHAPE.CUBE,
        'color': RGBColor(102, 102, 102),
        'icon_type': 'server',
        'width_ratio': 1.0,
        'height_ratio': 0.7,
    },
    'database': {
        'shape': MSO_SHAPE.CAN,
        'color': RGBColor(0, 153, 76),
        'icon_type': 'db',
        'width_ratio': 0.9,
        'height_ratio': 1.0,
    },
    'storage': {
        'shape': MSO_SHAPE.CAN,
        'color': RGBColor(0, 102, 153),
        'icon_type': 'storage',
        'width_ratio': 1.2,
        'height_ratio': 0.8,
    },
    
    # End Devices
    'pc': {
        'shape': MSO_SHAPE.RECTANGLE,
        'color': RGBColor(102, 153, 204),
        'icon_type': 'pc',
        'width_ratio': 1.0,
        'height_ratio': 0.8,
    },
    'laptop': {
        'shape': MSO_SHAPE.RECTANGLE,
        'color': RGBColor(102, 153, 204),
        'icon_type': 'laptop',
        'width_ratio': 1.1,
        'height_ratio': 0.7,
    },
    'phone': {
        'shape': MSO_SHAPE.ROUNDED_RECTANGLE,
        'color': RGBColor(153, 102, 204),
        'icon_type': 'phone',
        'width_ratio': 0.5,
        'height_ratio': 1.0,
    },
    
    # Cloud & External
    'cloud': {
        'shape': MSO_SHAPE.CLOUD,
        'color': RGBColor(66, 133, 244),
        'icon_type': 'cloud',
        'width_ratio': 1.3,
        'height_ratio': 0.9,
    },
    'internet': {
        'shape': MSO_SHAPE.OVAL,  # Use oval with globe-like styling
        'color': RGBColor(66, 133, 244),
        'icon_type': 'internet',
        'width_ratio': 1.0,
        'height_ratio': 1.0,
    },
    'datacenter': {
        'shape': MSO_SHAPE.RECTANGLE,
        'color': RGBColor(102, 102, 102),
        'icon_type': 'dc',
        'width_ratio': 1.2,
        'height_ratio': 1.0,
    },
}

# ============================================================================
# Network Zone Templates
# ============================================================================

ZONE_TEMPLATES = {
    'internet': {
        'title': '互联网接入区',
        'bg_color': RGBColor(245, 250, 255),
        'border_color': RGBColor(66, 133, 244),
        'border_width': 2,
    },
    'dmz': {
        'title': 'DMZ 区',
        'bg_color': RGBColor(255, 250, 245),
        'border_color': RGBColor(255, 153, 0),
        'border_width': 2,
    },
    'office': {
        'title': '办公运维区',
        'bg_color': RGBColor(245, 255, 245),
        'border_color': RGBColor(0, 153, 76),
        'border_width': 2,
    },
    'datacenter': {
        'title': '数据中心区',
        'bg_color': RGBColor(250, 245, 255),
        'border_color': RGBColor(153, 0, 204),
        'border_width': 2,
    },
    'security': {
        'title': '安全管理中心',
        'bg_color': RGBColor(255, 245, 245),
        'border_color': RGBColor(204, 0, 0),
        'border_width': 2,
    },
    'core': {
        'title': '核心交换区',
        'bg_color': RGBColor(245, 250, 255),
        'border_color': RGBColor(0, 102, 204),
        'border_width': 2,
    },
}

# ============================================================================
# Network Icon Drawer
# ============================================================================

class NetworkIconDrawer:
    """Draw network device icons with Cisco-style appearance"""
    
    def __init__(self, slide):
        self.slide = slide
    
    def draw_icon(self, icon_type, x, y, width, height, label=None, 
                  color=None, custom_shape=None):
        """Draw a network device icon"""
        
        icon_spec = NETWORK_ICONS.get(icon_type, NETWORK_ICONS['server'])
        shape_type = custom_shape or icon_spec['shape']
        icon_color = color or icon_spec['color']
        
        # Adjust dimensions based on icon type
        w = width * icon_spec['width_ratio']
        h = height * icon_spec['height_ratio']
        
        # Center the icon
        x_adj = x + (width - w) / 2
        y_adj = y + (height - h) / 2
        
        # Draw main shape
        shape = self.slide.shapes.add_shape(
            shape_type, Inches(x_adj), Inches(y_adj), Inches(w), Inches(h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = icon_color
        shape.line.color.rgb = RGBColor(50, 50, 50)
        shape.line.width = Pt(1.5)
        
        # Add icon-specific details
        self._add_icon_details(icon_type, shape, x_adj, y_adj, w, h)
        
        # Add label
        if label:
            self._add_label(x, y, width, height, label)
        
        return shape
    
    def _add_icon_details(self, icon_type, shape, x, y, w, h):
        """Add icon-specific visual details"""
        
        if icon_type in ['router', 'switch']:
            # Add arrows inside oval
            self._add_arrows(shape, x, y, w, h, icon_type)
        elif icon_type == 'firewall':
            # Add brick pattern hint
            self._add_brick_pattern(shape, x, y, w, h)
        elif icon_type == 'server':
            # Add server slots
            self._add_server_slots(shape, x, y, w, h)
        elif icon_type == 'pc':
            # Add monitor stand
            self._add_monitor_stand(shape, x, y, w, h)
        elif icon_type == 'internet':
            # Add globe-like pattern
            self._add_globe_pattern(shape, x, y, w, h)
    
    def _add_arrows(self, shape, x, y, w, h, icon_type):
        """Add arrow indicators for router/switch"""
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        if icon_type == 'router':
            p.text = '⇄'
        elif icon_type == 'switch':
            p.text = '⇅'
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True
    
    def _add_brick_pattern(self, shape, x, y, w, h):
        """Add firewall brick pattern hint"""
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = ''
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    
    def _add_server_slots(self, shape, x, y, w, h):
        """Add server rack slots"""
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = '≡'
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.alignment = PP_ALIGN.CENTER
    
    def _add_monitor_stand(self, shape, x, y, w, h):
        """Add PC monitor stand"""
        pass
    
    def _add_globe_pattern(self, shape, x, y, w, h):
        """Add globe-like pattern for internet icon"""
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = '🌐'
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    
    def _add_label(self, x, y, w, h, label):
        """Add label below icon"""
        label_box = self.slide.shapes.add_textbox(
            Inches(x), Inches(y + h), Inches(w), Inches(0.4)
        )
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.alignment = PP_ALIGN.CENTER
        label_box.fill.background()
        label_box.line.fill.background()

# ============================================================================
# Network Zone Container
# ============================================================================

class NetworkZone:
    """Create network zone containers"""
    
    def __init__(self, slide, zone_type='custom'):
        self.slide = slide
        self.zone_type = zone_type
        self.template = ZONE_TEMPLATES.get(zone_type, {
            'title': '网络区域',
            'bg_color': RGBColor(245, 245, 245),
            'border_color': RGBColor(150, 150, 150),
            'border_width': 2,
        })
    
    def create_zone(self, x, y, width, height, title=None, devices=None):
        """Create a network zone container"""
        
        # Zone background
        zone_bg = self.slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y), Inches(width), Inches(height)
        )
        zone_bg.fill.solid()
        zone_bg.fill.fore_color.rgb = self.template['bg_color']
        zone_bg.line.color.rgb = self.template['border_color']
        zone_bg.line.width = Pt(self.template['border_width'])
        
        # Zone title
        zone_title = title or self.template['title']
        title_box = self.slide.shapes.add_textbox(
            Inches(x + width - 1.5), Inches(y + 0.1), Inches(1.3), Inches(0.3)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = zone_title
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.alignment = PP_ALIGN.RIGHT
        title_box.fill.background()
        title_box.line.fill.background()
        
        return zone_bg

# ============================================================================
# Network Topology Builder
# ============================================================================

class NetworkTopologyBuilder:
    """Build complete network topology diagrams"""
    
    def __init__(self, slide):
        self.slide = slide
        self.icon_drawer = NetworkIconDrawer(slide)
        self.devices = {}  # Store device positions for connections
    
    def add_zone(self, zone_id, zone_type, x, y, width, height, title=None):
        """Add a network zone"""
        zone = NetworkZone(self.slide, zone_type)
        zone.create_zone(x, y, width, height, title)
        return zone
    
    def add_device(self, device_id, device_type, x, y, width=0.8, height=0.6, 
                   label=None, color=None, zone_id=None):
        """Add a network device"""
        shape = self.icon_drawer.draw_icon(
            device_type, x, y, width, height, label, color
        )
        
        # Store position for connections
        self.devices[device_id] = {
            'x': x + width / 2,
            'y': y + height / 2,
            'width': width,
            'height': height,
            'shape': shape,
        }
        
        return shape
    
    def add_connection(self, device1_id, device2_id, label=None, 
                       color=RGBColor(100, 100, 100), width=1.5, dashed=False):
        """Add a connection between two devices"""
        if device1_id not in self.devices or device2_id not in self.devices:
            return None
        
        dev1 = self.devices[device1_id]
        dev2 = self.devices[device2_id]
        
        # Create connector
        connector = self.slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(dev1['x']), Inches(dev1['y']),
            Inches(dev2['x']), Inches(dev2['y'])
        )
        connector.line.color.rgb = color
        connector.line.width = Pt(width)
        if dashed:
            connector.line.dash_style = 4
        
        # Add label if provided
        if label:
            mid_x = (dev1['x'] + dev2['x']) / 2
            mid_y = (dev1['y'] + dev2['y']) / 2
            
            label_box = self.slide.shapes.add_textbox(
                Inches(mid_x - 0.4), Inches(mid_y - 0.15),
                Inches(0.8), Inches(0.3)
            )
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = label
            p.font.size = Pt(8)
            p.font.color.rgb = RGBColor(80, 80, 80)
            p.alignment = PP_ALIGN.CENTER
            label_box.fill.solid()
            label_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
            label_box.line.fill.background()
        
        return connector

# ============================================================================
# Physical Deployment Diagram Builder (参考图)
# ============================================================================

class PhysicalDeploymentBuilder:
    """Build physical deployment diagrams like the reference image"""
    
    def __init__(self, slide):
        self.slide = slide
        self.topology = NetworkTopologyBuilder(slide)
    
    def build_reference_diagram(self):
        """Build the reference physical deployment diagram"""
        slide_w = 13.333
        slide_h = 7.5
        
        # ========== Create Zones ==========
        
        # Internet Access Zone (Top Center)
        self.topology.add_zone(
            'internet_zone', 'internet',
            4.5, 0.8, 4.0, 2.2
        )
        
        # Core Zone (Center)
        self.topology.add_zone(
            'core_zone', 'core',
            4.5, 3.2, 4.0, 1.5
        )
        
        # Office Zone (Bottom Center)
        self.topology.add_zone(
            'office_zone', 'office',
            4.5, 4.9, 4.0, 2.4
        )
        
        # DMZ Zone (Right)
        self.topology.add_zone(
            'dmz_zone', 'dmz',
            9.2, 1.5, 4.0, 4.5
        )
        
        # Security Management Center (Left)
        self.topology.add_zone(
            'security_zone', 'security',
            0.2, 2.5, 4.0, 3.0
        )
        
        # ========== Add Devices ==========
        
        # Internet Cloud
        self.topology.add_device(
            'internet', 'internet',
            6.2, 1.0, 0.7, 0.7,
            label='Internet'
        )
        
        # Border Firewalls (Internet Zone)
        self.topology.add_device(
            'firewall1', 'firewall',
            5.0, 2.0, 0.6, 0.6,
            label='边界防火墙'
        )
        self.topology.add_device(
            'firewall2', 'firewall',
            7.4, 2.0, 0.6, 0.6,
            label='边界防火墙'
        )
        
        # Core Switches
        self.topology.add_device(
            'core_switch1', 'switch',
            5.2, 3.6, 0.7, 0.5,
            label='核心交换机'
        )
        self.topology.add_device(
            'core_switch2', 'switch',
            7.0, 3.6, 0.7, 0.5,
            label='核心交换机'
        )
        
        # Access Switch (Office Zone)
        self.topology.add_device(
            'access_switch', 'switch',
            6.2, 5.3, 0.7, 0.5,
            label='接入交换机'
        )
        
        # End Devices (PCs)
        for i in range(5):
            self.topology.add_device(
                f'pc_{i}', 'pc',
                4.8 + i * 0.7, 6.2, 0.5, 0.4,
                label='终端'
            )
        
        # DMZ Devices
        # Firewall chain
        self.topology.add_device(
            'dmz_firewall', 'firewall',
            8.5, 3.5, 0.6, 0.6,
            label='边界防火墙'
        )
        self.topology.add_device(
            'web_firewall', 'firewall',
            9.3, 3.5, 0.6, 0.6,
            label='WEB 应用防火墙'
        )
        self.topology.add_device(
            'dmz_switch', 'switch',
            10.2, 3.5, 0.7, 0.5,
            label='接入交换机'
        )
        
        # Application Servers
        self.topology.add_device(
            'app_server1', 'server',
            10.0, 1.8, 0.5, 0.35,
            label='应用服务器'
        )
        self.topology.add_device(
            'app_server2', 'server',
            10.7, 1.8, 0.5, 0.35,
            label='应用服务器'
        )
        self.topology.add_device(
            'app_server3', 'server',
            11.4, 1.8, 0.5, 0.35,
            label='应用服务器'
        )
        
        # Database Servers
        self.topology.add_device(
            'db_server1', 'server',
            10.0, 2.7, 0.5, 0.35,
            label='数据库服务器'
        )
        self.topology.add_device(
            'db_server2', 'server',
            10.7, 2.7, 0.5, 0.35,
            label='数据库服务器'
        )
        self.topology.add_device(
            'db_server3', 'server',
            11.4, 2.7, 0.5, 0.35,
            label='数据库服务器'
        )
        
        # Public Servers
        self.topology.add_device(
            'public_server1', 'server',
            10.0, 4.8, 0.5, 0.35,
            label='公共服务器'
        )
        self.topology.add_device(
            'public_server2', 'server',
            10.7, 4.8, 0.5, 0.35,
            label='公共服务器'
        )
        self.topology.add_device(
            'public_server3', 'server',
            11.4, 4.8, 0.5, 0.35,
            label='公共服务器'
        )
        
        # Security Management Center Devices
        self.topology.add_device(
            'sec_switch', 'switch',
            3.0, 3.5, 0.7, 0.5,
            label='接入交换机'
        )
        self.topology.add_device(
            'sec_firewall', 'firewall',
            3.9, 3.5, 0.6, 0.6,
            label='边界防火墙'
        )
        
        # Security Systems
        sec_systems = [
            ('log_audit', '日志审计系统', 1.0, 2.8),
            ('db_audit', '数据库审计系统', 1.0, 3.4),
            ('ops_audit', '运维审计系统', 1.0, 4.0),
            ('net_mgmt', '网络管理系统', 1.0, 4.6),
            ('security_ctrl', '杀毒软件及统一管控系统', 1.0, 5.2),
        ]
        
        for dev_id, label, x, y in sec_systems:
            self.topology.add_device(
                dev_id, 'server',
                x, y, 0.6, 0.4,
                label=label
            )
        
        # ========== Add Connections ==========
        
        # Internet to Firewalls
        self.topology.add_connection('internet', 'firewall1')
        self.topology.add_connection('internet', 'firewall2')
        
        # Firewalls to Core Switches
        self.topology.add_connection('firewall1', 'core_switch1')
        self.topology.add_connection('firewall1', 'core_switch2')
        self.topology.add_connection('firewall2', 'core_switch1')
        self.topology.add_connection('firewall2', 'core_switch2')
        
        # Core Switches to Access Switch
        self.topology.add_connection('core_switch1', 'access_switch')
        self.topology.add_connection('core_switch2', 'access_switch')
        
        # Access Switch to PCs
        for i in range(5):
            self.topology.add_connection('access_switch', f'pc_{i}')
        
        # Core to DMZ
        self.topology.add_connection('core_switch1', 'dmz_firewall')
        self.topology.add_connection('dmz_firewall', 'web_firewall')
        self.topology.add_connection('web_firewall', 'dmz_switch')
        
        # DMZ Switch to Servers
        for i in range(1, 4):
            self.topology.add_connection('dmz_switch', f'app_server{i}')
            self.topology.add_connection('dmz_switch', f'db_server{i}')
            self.topology.add_connection('dmz_switch', f'public_server{i}')
        
        # Security Center
        self.topology.add_connection('core_switch1', 'sec_firewall')
        self.topology.add_connection('sec_firewall', 'sec_switch')
        
        for dev_id, _, _, _ in sec_systems:
            self.topology.add_connection('sec_switch', dev_id)
        
        # Add title
        title_box = self.slide.shapes.add_textbox(
            Inches(0.5), Inches(0.1), Inches(12), Inches(0.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = '物理环境部署图'
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER
    
    def build_custom_topology(self, topology_config):
        """Build custom topology from configuration"""
        
        # Create zones
        for zone in topology_config.get('zones', []):
            zone_type = zone.get('type', 'custom')
            self.topology.add_zone(
                zone['id'], zone_type,
                zone['x'], zone['y'], zone['width'], zone['height'],
                zone.get('title')
            )
        
        # Create devices
        for device in topology_config.get('devices', []):
            self.topology.add_device(
                device['id'],
                device.get('type', 'server'),
                device['x'], device['y'],
                device.get('width', 0.8),
                device.get('height', 0.6),
                device.get('label'),
                device.get('color'),
                device.get('zone_id')
            )
        
        # Create connections
        for conn in topology_config.get('connections', []):
            self.topology.add_connection(
                conn['from'],
                conn['to'],
                conn.get('label'),
                conn.get('color', RGBColor(100, 100, 100)),
                conn.get('width', 1.5),
                conn.get('dashed', False)
            )

# ============================================================================
# Main Generation Functions
# ============================================================================

def create_physical_deployment_diagram(output_path, title='物理环境部署图'):
    """Create the reference physical deployment diagram"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Create blank slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # White background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.fill.background()
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.1), Inches(12), Inches(0.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER
    
    # Build diagram
    builder = PhysicalDeploymentBuilder(slide)
    builder.build_reference_diagram()
    
    prs.save(output_path)
    print(f"✓ Physical deployment diagram generated: {output_path}")

def create_custom_network_topology(output_path, topology_config, title='网络拓扑图'):
    """Create custom network topology from configuration"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # White background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.fill.background()
    
    # Add title
    if title:
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.1), Inches(12), Inches(0.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER
    
    # Build diagram
    builder = PhysicalDeploymentBuilder(slide)
    builder.build_custom_topology(topology_config)
    
    prs.save(output_path)
    print(f"✓ Network topology diagram generated: {output_path}")

# ============================================================================
# CLI
# ============================================================================

def main():
    parser = argparse.ArgumentParser(description='Network Topology Diagram Generator')
    parser.add_argument('--input', '-i', required=True, help='Input JSON config file')
    parser.add_argument('--output', '-o', required=True, help='Output PPTX file')
    parser.add_argument('--title', '-t', default='网络拓扑图', help='Diagram title')
    parser.add_argument('--template', choices=['reference', 'custom'], default='custom',
                       help='Use reference diagram template or custom config')
    
    args = parser.parse_args()
    
    # Load config
    with open(args.input, 'r', encoding='utf-8') as f:
        config = json.load(f)
    
    if args.template == 'reference':
        create_physical_deployment_diagram(args.output, args.title)
    else:
        create_custom_network_topology(args.output, config, args.title)

if __name__ == '__main__':
    main()
